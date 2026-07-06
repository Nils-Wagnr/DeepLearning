"""Create a fully German ClaimGuard deck from the editable English presentation."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation


HERE = Path(__file__).resolve().parent
SOURCE = HERE / "ClaimGuard_Capstone_Presentation.pptx"
OUTPUT = HERE / "ClaimGuard_Capstone_Praesentation_DE.pptx"
NOTES_OUTPUT = HERE / "ClaimGuard_Praesentation_Notizen_DE.md"


TRANSLATIONS = {
    "15 MIN • TEAM PRESENTATION": "15 MIN • TEAMPRÄSENTATION",
    "Evidence-backed academic integrity checking": "Evidenzbasierte Prüfung akademischer Integrität",
    "From PDF claims to reference validation and source-support verdicts": "Von Aussagen im PDF zur Referenzprüfung und Quellenbewertung",
    "MWI Deep Learning Capstone  •  2 July 2026": "MWI Deep Learning Capstone  •  3. Juli 2026",
    "PROBLEM MOTIVATION": "PROBLEMMOTIVATION",
    "Why academic integrity checking matters": "Warum die Prüfung akademischer Integrität wichtig ist",
    "A citation can be real…": "Eine Quelle kann echt sein …",
    "…and still not support the sentence that cites it.": "… und trotzdem die zitierte Aussage nicht stützen.",
    "Manual review is slow; naive automation creates false accusations.": "Manuelle Prüfung ist langsam; naive Automatisierung erzeugt falsche Verdächtigungen.",
    "TRIAGE, NOT A VERDICT": "VORPRÜFUNG, KEIN URTEIL",
    "•  Scale  Long reports contain dozens of checkable claims.": "•  Umfang  Lange Berichte enthalten viele überprüfbare Aussagen.",
    "•  Subtlety  Existence, relevance, and factual support are different questions.": "•  Feinheit  Existenz, Relevanz und faktische Stützung sind verschiedene Fragen.",
    "•  Risk  False positives can unfairly damage students or researchers.": "•  Risiko  Fehlalarme können Studierenden und Forschenden unfair schaden.",
    "•  Need  Reviewers require traceable evidence, not a black-box score.": "•  Bedarf  Prüfende brauchen nachvollziehbare Evidenz statt einer Blackbox-Zahl.",
    "Goal: make suspicious passages faster to inspect and easier to explain.": "Ziel: Auffällige Passagen schneller prüfen und verständlich begründen.",
    "PROBLEM FRAMING": "PROBLEMSTELLUNG",
    "Three checks — one auditable report": "Drei Prüfungen — ein nachvollziehbarer Bericht",
    "Does this need a citation?": "Benötigt das eine Quelle?",
    "Claim type, section, citation marker, severity": "Aussagetyp, Abschnitt, Zitiermarker, Schweregrad",
    "Does the source exist?": "Existiert die Quelle?",
    "IEEE parsing plus scholarly metadata APIs": "IEEE-Parsing plus wissenschaftliche Metadaten-APIs",
    "Does the source support it?": "Stützt die Quelle die Aussage?",
    "Evidence retrieval plus five-way verification": "Evidenzsuche plus fünfstufige Verifikation",
    "Output: human-readable Markdown + complete JSON provenance for evaluation and audit": "Ausgabe: lesbares Markdown + vollständige JSON-Provenienz für Evaluation und Audit",
    "SYSTEM ARCHITECTURE": "SYSTEMARCHITEKTUR",
    "Pipeline architecture": "Pipeline-Architektur",
    "Ingest": "Einlesen",
    "PDF/TXT extraction\vsection cleanup": "PDF/TXT-Extraktion\vAbschnittsbereinigung",
    "Claims": "Aussagen",
    "type + citation need\vcontext linking": "Typ + Quellenbedarf\vKontextverknüpfung",
    "References": "Referenzen",
    "IEEE parsing\vmetadata validation": "IEEE-Parsing\vMetadatenprüfung",
    "Evidence": "Evidenz",
    "lexical / embeddings\vchunk retrieval": "lexikalisch / Embeddings\vChunk-Suche",
    "Verdict": "Bewertung",
    "heuristic / local /\vfrontier / LoRA": "Heuristik / lokal /\vFrontier / LoRA",
    "AUDIT RAIL": "AUDIT-SPUR",
    "evidence • confidence • model ID • latency • token usage • response/request IDs": "Evidenz • Konfidenz • Modell-ID • Latenz • Tokenverbrauch • Antwort-/Anfrage-IDs",
    "Offline first": "Offline zuerst",
    "Deterministic fallback": "Regelbasierter Fallback",
    "Network features are opt-in": "Netzwerkfunktionen nur optional",
    "Same schema": "Ein Schema",
    "Comparable backends": "Vergleichbare Backends",
    "Five support labels": "Fünf Unterstützungslabels",
    "Human review": "Mensch prüft",
    "Final decision": "Endgültige Entscheidung",
    "No automatic misconduct claim": "Kein automatischer Täuschungsvorwurf",
    "TECHNICAL DEPTH": "TECHNISCHE TIEFE",
    "Technical decisions and trade-offs": "Technische Entscheidungen und Zielkonflikte",
    "RETRIEVAL": "EVIDENZSUCHE",
    "Semantic matching; cosine fallback on Windows": "Semantischer Abgleich; Kosinus-Fallback unter Windows",
    "METADATA": "METADATEN",
    "Coverage across DOI, title, authors, OA status": "Abdeckung von DOI, Titel, Autoren und Open-Access-Status",
    "LOCAL MODEL": "LOKAL",
    "Private and offline; slower on this machine": "Privat und offline; auf diesem Rechner langsamer",
    "FRONTIER": "FRONTIER",
    "Fast structured output; paid and external": "Schnelle strukturierte Ausgabe; kostenpflichtig und extern",
    "FINE-TUNED": "FINE-TUNING",
    "Cheap local inference; three-label taxonomy": "Günstige lokale Inferenz; Taxonomie mit drei Labels",
    "Design principle: graceful degradation — APIs or embeddings may fail without stopping the report.": "Designprinzip: kontrollierte Degradation — API- oder Embedding-Fehler stoppen den Bericht nicht.",
    "Demo on our real RQ5 report": "Demo mit unserem echten RQ5-Bericht",
    "sentences": "Sätze",
    "factual claims": "faktische Aussagen",
    "citation flags": "Quellenhinweise",
    "REAL CLAIM FROM SENTENCE 11": "ECHTE AUSSAGE AUS SATZ 11 (ÜBERSETZT)",
    "“RAG connects a language model to external knowledge without changing model weights.”": "„RAG verbindet ein Sprachmodell mit externem Wissen, ohne Modellgewichte zu verändern.“",
    "Citations detected: [1, 3]  •  Top evidence retained": "Erkannte Quellen: [1, 3]  •  Beste Evidenz gespeichert",
    "insufficient\vevidence": "unzureichende\vEvidenz",
    "supported": "gestützt",
    "40.902 s • local": "40,902 s • lokal",
    "Demo message: the evidence is visible — disagreement becomes reviewable instead of hidden.": "Kernaussage: Die Evidenz ist sichtbar — Unterschiede werden prüfbar statt verborgen.",
    "EVALUATION": "EVALUATION",
    "Key results: controlled tests vs. real reports": "Zentrale Ergebnisse: kontrollierte Tests vs. echte Berichte",
    "IEEE parser accuracy": "Genauigkeit des IEEE-Parsers",
    "5-field-complete test cases": "Fünf Testfälle mit vollständigen Feldern",
    "Citation-needed F1": "F1 für Quellenbedarf",
    "50 manually labeled RQ sentences": "50 manuell annotierte RQ-Sätze",
    "SciFact validation Macro-F1": "SciFact-Validierung: Macro-F1",
    "450 validation pairs": "450 Validierungspaare",
    "Main finding: synthetic benchmarks overestimate how well simple rules transfer to authentic PDFs.": "Hauptergebnis: Synthetische Benchmarks überschätzen den Transfer einfacher Regeln auf echte PDFs.",
    "LOCAL VS. FRONTIER": "LOKAL VS. FRONTIER",
    "Model comparison: speed, control, and taxonomy": "Modellvergleich: Geschwindigkeit, Kontrolle und Taxonomie",
    "Backend": "Backend",
    "Observed result": "Beobachtetes Ergebnis",
    "Latency": "Latenz",
    "Best use": "Geeigneter Einsatz",
    "Limitation": "Einschränkung",
    "Heuristic": "Heuristik",
    "Transparent baseline": "Transparente Baseline",
    "Relies on surface patterns": "Beruht auf Oberflächenmustern",
    "Cheap local inference": "Günstige lokale Inferenz",
    "SciFact has 3 vs. 5 labels": "SciFact hat 3 statt 5 Labels",
    "supported†": "gestützt†",
    "Private local review": "Private lokale Prüfung",
    "Slow; one diagnostic claim": "Langsam; eine diagnostische Aussage",
    "insufficient†": "unzureichend†",
    "Structured frontier review": "Strukturierte Frontier-Prüfung",
    "Paid/external; one diagnostic claim": "Kostenpflichtig/extern; eine diagnostische Aussage",
    "* 30-case gold benchmark     † same real RQ5 claim; diagnostic only — not an accuracy comparison": "* Gold-Benchmark mit 30 Fällen     † dieselbe echte RQ5-Aussage; nur diagnostisch — kein Genauigkeitsvergleich",
    "LIMITS": "GRENZEN",
    "What works — and where it breaks": "Was funktioniert — und wo es scheitert",
    "WORKS WELL": "FUNKTIONIERT GUT",
    "•  End-to-end PDF → Markdown/JSON workflow": "•  Durchgängiger Ablauf von PDF zu Markdown/JSON",
    "•  Traceable evidence and backend metadata": "•  Nachvollziehbare Evidenz und Backend-Metadaten",
    "•  Robust IEEE parsing for Report 5": "•  Robustes IEEE-Parsing für Bericht 5",
    "•  Offline fallback and real CUDA LoRA path": "•  Offline-Fallback und echter CUDA-LoRA-Pfad",
    "LIMITATIONS": "EINSCHRÄNKUNGEN",
    "•  PDF extraction introduces spacing/heading noise": "•  PDF-Extraktion erzeugt Abstands- und Überschriftenfehler",
    "•  Real claim detection generalizes poorly": "•  Erkennung echter Aussagen generalisiert schwach",
    "•  API metadata may time out or rate-limit": "•  Metadaten-APIs können ausfallen oder drosseln",
    "•  Binoculars' calibrated 14B pair exceeds 8 GB VRAM": "•  Binoculars' kalibriertes 14B-Paar übersteigt 8 GB VRAM",
    "Ethical boundary: a flag prioritizes human review; it never proves misconduct or AI authorship.": "Ethische Grenze: Ein Hinweis priorisiert menschliche Prüfung; er beweist weder Täuschung noch KI-Autorschaft.",
    "CONCLUSION": "FAZIT",
    "Takeaways": "Kernaussagen",
    "ClaimGuard makes integrity review auditable.": "ClaimGuard macht Integritätsprüfungen nachvollziehbar.",
    "Pipeline": "Pipeline",
    "Claims, references, evidence, and verdicts in one workflow.": "Aussagen, Referenzen, Evidenz und Bewertungen in einem Ablauf.",
    "Real-report evaluation exposes where synthetic scores mislead.": "Evaluation mit echten Berichten zeigt, wo synthetische Werte täuschen.",
    "Responsibility": "Verantwortung",
    "Use automation to focus human attention — not replace judgment.": "Automatisierung lenkt menschliche Aufmerksamkeit — sie ersetzt kein Urteil.",
    "Next: expand double-annotated real claims and run the full OpenAI/Ollama gold benchmark.": "Nächster Schritt: doppelt annotierte Echtdaten erweitern und den vollständigen OpenAI/Ollama-Benchmark ausführen.",
    "Nick closes • Nils transitions to Q&A": "Nick fasst zusammen • Nils leitet zur Fragerunde über",
    "Questions & discussion": "Fragen & Diskussion",
    "5 minutes": "5 Minuten",
    "Architecture  •  Evaluation  •  Model trade-offs  •  Ethics": "Architektur  •  Evaluation  •  Modell-Zielkonflikte  •  Ethik",
    "05:00 Q&A": "05:00 FRAGEN",
}


CHART_TRANSLATIONS = {
    "Real RQ\nclaims": "Echte RQ-\nAussagen",
    "Lexical\nverifier": "Lexikalische\nPrüfung",
    "Embedding\nverifier": "Embedding-\nPrüfung",
    "LoRA\ntransfer": "LoRA-\nTransfer",
    "Macro-F1": "Macro-F1",
}


GERMAN_NOTES = [
    "Nils — 20 Sekunden\nMit einem Satz eröffnen: ClaimGuard entscheidet nicht über Täuschung, sondern ordnet Evidenz für eine menschliche Prüfung. Beide Vortragenden sowie die drei Fragen Aussage, Referenz und Stützung vorstellen.",
    "Nils — 55 Sekunden\nDen Kontrast links nutzen. Plagiat, erfundene Referenzen und nicht gestützte Zitate sind unterschiedliche Probleme. ClaimGuard konzentriert sich auf Zitierintegrität und Evidenz. Das System unterstützt Entscheidungen; es erhebt keinen Vorwurf.",
    "Nils — 45 Sekunden\nVon links nach rechts erklären. Die drei Fragen entsprechen den drei Kernmodulen. Jede Entscheidung bleibt im JSON nachvollziehbar; Markdown ist die Ansicht für Prüfende.",
    "Nils — 70 Sekunden\nDen fünf nummerierten Schritten folgen. Referenzen begrenzen die Evidenzsuche; der Verifikator erhält nicht einfach das gesamte Internet. Die Audit-Spur macht die Ausgabe reproduzierbar. Danach an Nick übergeben.",
    "Nick — 75 Sekunden\nEs gibt kein einzelnes bestes Backend. Die APIs prüfen bibliografische Daten, nicht die inhaltliche Stützung. Die Evidenzsuche nutzt Embeddings und fällt bei Bedarf auf lexikalische Suche zurück. OpenAI nutzt store=false; lokale Modelle halten Texte auf dem Gerät.",
    "Nick — 45 Sekunden\nDas eingebettete UI-Video startet automatisch. Zuerst wird der echte RQ5-Bericht hochgeladen und lokal analysiert. Im Überblick nur kurz auf Prüfbedarf und die methodenunabhängigen Dokumentkennzahlen zeigen. Danach wechselt die Demo in den Modellvergleich: Heuristik und Ollama erhalten dieselben Claims und dieselbe Evidenz. Abweichungen sind kein automatischer Fehlernachweis, sondern gezielte Stellen für die menschliche Prüfung. Falls das Video nicht startet, ClaimGuard_UI_Demo.mp4 separat öffnen oder die statische Originalfolie verwenden.",
    "Nick — 75 Sekunden\nMit dem F1-Wert 0,900 der lexikalischen Prüfung beginnen und mit 0,437 Macro-F1 bei echten Aussagen kontrastieren. Embeddings halfen in diesem kleinen Direkt-Evidenz-Benchmark nicht. Das Parser-Ergebnis vorsichtig einordnen: nur fünf Fälle. Entscheidend ist die Lücke zwischen synthetischen und echten Daten.",
    "Nils — 75 Sekunden\nGold-Benchmark-Ergebnisse von der mit † markierten Einzelfall-Diagnose trennen. Die Heuristik gewinnt den kleinen synthetischen Benchmark, weil direkte lexikalische Muster belohnt werden. LoRA ist schnell, kann aber zwei der fünf ClaimGuard-Labels nicht ausgeben. OpenAI war deutlich schneller als Ollama; ein vollständiger Goldvergleich fehlt noch.",
    "Nils — 55 Sekunden\nJe zwei konkrete Stärken und Schwächen nennen. Einschränkungen sind gemessene technische Erkenntnisse. Binoculars ist integriert, aber das kalibrierte Modellpaar ist mit 8 GB VRAM nicht ausführbar.",
    "Nick — 20 Sekunden; Nils — 15 Sekunden\nNick fasst Pipeline und Evidenz zusammen. Nils schließt mit der ethischen Grenze und lädt zu Fragen ein.",
    "Beide — 5 Minuten\nVorgeschlagene Aufteilung: Nils beantwortet Fragen zu Architektur und Evaluation; Nick übernimmt Modelle, APIs und Demo. Bei Fragen zum niedrigen Echtdatenwert zuerst auf Domain Shift und die Einzelannotation eingehen.",
]


def replace_paragraph_text(paragraph, new_text: str) -> None:
    """Replace text while keeping the first run's formatting."""
    if paragraph.runs:
        parts = new_text.split("\v")
        paragraph.runs[0].text = parts[0]
        for run in paragraph.runs[1:]:
            run.text = ""
        for part in parts[1:]:
            paragraph.add_line_break()
            paragraph.add_run().text = part
    else:
        paragraph.text = new_text.replace("\v", "\n")


def translate_visible_text(prs: Presentation) -> set[str]:
    untranslated: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                old = paragraph.text
                if old in TRANSLATIONS:
                    replace_paragraph_text(paragraph, TRANSLATIONS[old])
                elif old.strip() and any(ch.isalpha() for ch in old):
                    untranslated.add(old)
    return untranslated


def translate_chart(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_chart", False):
                continue
            for value in shape.chart._chartSpace.xpath(".//c:v"):
                if value.text in CHART_TRANSLATIONS:
                    value.text = CHART_TRANSLATIONS[value.text]


def write_notes(prs: Presentation) -> None:
    lines = ["# ClaimGuard – Präsentationsnotizen", "", "10 Minuten Vortrag + 5 Minuten Fragen.", ""]
    for index, (slide, note) in enumerate(zip(prs.slides, GERMAN_NOTES, strict=True), start=1):
        slide.notes_slide.notes_text_frame.text = note
        lines.extend([f"## Folie {index}", "", note, ""])
    NOTES_OUTPUT.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)
    prs = Presentation(SOURCE)
    if len(prs.slides) != len(GERMAN_NOTES):
        raise RuntimeError(f"Erwartet: {len(GERMAN_NOTES)} Folien; gefunden: {len(prs.slides)}")
    untranslated = translate_visible_text(prs)
    translate_chart(prs)
    write_notes(prs)
    prs.core_properties.title = "ClaimGuard — Evidenzbasierte Prüfung akademischer Integrität"
    prs.core_properties.subject = "Deutsche MWI-Deep-Learning-Capstone-Präsentation"
    prs.save(OUTPUT)
    print(f"Erstellt: {OUTPUT}")
    print(f"Erstellt: {NOTES_OUTPUT}")
    print("Bewusst unverändert (Namen, Befehle, Modelle und Messwerte):")
    for text in sorted(untranslated):
        print(f"  {text!r}")


if __name__ == "__main__":
    main()
