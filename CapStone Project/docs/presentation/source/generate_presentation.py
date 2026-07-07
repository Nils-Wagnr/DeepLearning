"""Generate the ClaimGuard capstone presentation as an editable PowerPoint deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
BUILD_DIR = HERE / "build"
OUT = BUILD_DIR / "ClaimGuard_Capstone_Presentation.pptx"
NOTES_OUT = BUILD_DIR / "ClaimGuard_Presentation_Notes.md"

NAVY = "0B132B"
NAVY_2 = "16213E"
BLUE = "2563EB"
CYAN = "38BDF8"
GREEN = "10B981"
AMBER = "F59E0B"
RED = "EF4444"
PURPLE = "8B5CF6"
WHITE = "FFFFFF"
INK = "182230"
MUTED = "5D6B7A"
LIGHT = "F3F6FA"
LINE = "D7DFEA"

W = Inches(13.333)
H = Inches(7.5)


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_fill(shape, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_line(shape, color: str, width: float = 1.0) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)


def textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 20,
    color: str = INK,
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.05,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return shape


def box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = LINE,
    radius: bool = True,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line, 1)
    if radius and hasattr(shape, "adjustments") and len(shape.adjustments):
        shape.adjustments[0] = 0.12
    return shape


def pill(slide, x: float, y: float, w: float, text: str, fill: str, color: str = WHITE):
    shape = box(slide, x, y, w, 0.34, fill=fill, line=fill)
    textbox(slide, x, y + 0.01, w, 0.3, text, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
    return shape


def add_bullets(slide, x, y, w, h, items, *, size=18, color=INK, bullet_color=None, gap=8):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for index, item in enumerate(items):
        if isinstance(item, tuple):
            title, detail = item
            text = f"{title}  {detail}"
        else:
            title, detail, text = "", "", str(item)
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"•  {text}"
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        if title:
            p.runs[0].font.bold = True
    return shape


def add_header(slide, number: int, title: str, speaker: str, timing: str, section: str):
    textbox(slide, 0.55, 0.25, 8.8, 0.3, section.upper(), size=10, color=BLUE, bold=True)
    textbox(slide, 0.55, 0.57, 11.7, 0.55, title, size=27, color=INK, bold=True)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(0.55), Inches(7.08), Inches(12.78), Inches(7.08)
    )
    set_line(line, LINE, 1)
    textbox(slide, 0.58, 7.12, 4.0, 0.22, f"{speaker}  •  {timing}", size=9, color=MUTED)
    textbox(slide, 12.25, 7.12, 0.5, 0.22, f"{number:02d}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_notes(slide, notes: str) -> None:
    slide.notes_slide.notes_text_frame.text = notes


def new_slide(prs, number, title, speaker, timing, section):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fill(slide.background, WHITE)
    add_header(slide, number, title, speaker, timing, section)
    return slide


def metric_card(slide, x, y, w, value, label, color=BLUE, note=None):
    box(slide, x, y, w, 1.35, fill=LIGHT, line=LINE)
    textbox(slide, x + 0.18, y + 0.15, w - 0.36, 0.48, value, size=27, color=color, bold=True)
    textbox(slide, x + 0.18, y + 0.68, w - 0.36, 0.26, label, size=12, color=INK, bold=True)
    if note:
        textbox(slide, x + 0.18, y + 0.99, w - 0.36, 0.22, note, size=9, color=MUTED)


def add_architecture_stage(slide, x, y, w, index, title, detail, color):
    box(slide, x, y, w, 1.38, fill=WHITE, line=color)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.12), Inches(y + 0.15), Inches(0.42), Inches(0.42))
    set_fill(circle, color)
    set_line(circle, color)
    textbox(slide, x + 0.12, y + 0.17, 0.42, 0.34, str(index), size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, x + 0.66, y + 0.15, w - 0.78, 0.38, title, size=15, color=INK, bold=True)
    textbox(slide, x + 0.16, y + 0.67, w - 0.32, 0.52, detail, size=10.5, color=MUTED)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    prs.core_properties.title = "ClaimGuard — Academic Integrity Checking"
    prs.core_properties.subject = "MWI Deep Learning Capstone presentation"
    prs.core_properties.author = "Nils Wagner & Nick Wenzel"

    # 1 — Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_fill(s.background, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.14), H)
    set_fill(accent, CYAN)
    set_line(accent, CYAN)
    shield = s.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(0.82), Inches(0.8), Inches(0.76), Inches(0.76))
    set_fill(shield, BLUE)
    set_line(shield, CYAN, 1.5)
    textbox(s, 0.82, 1.0, 0.76, 0.3, "CG", size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    pill(s, 9.95, 0.72, 2.45, "15 MIN • TEAM PRESENTATION", BLUE)
    textbox(s, 0.82, 2.05, 11.2, 0.8, "ClaimGuard", size=44, color=WHITE, bold=True)
    textbox(s, 0.82, 2.88, 10.9, 0.8, "Evidence-backed academic integrity checking", size=25, color=CYAN, bold=True)
    textbox(s, 0.84, 3.88, 9.8, 0.72, "From PDF claims to reference validation and source-support verdicts", size=18, color="CBD5E1")
    textbox(s, 0.84, 5.52, 5.0, 0.3, "Nils Wagner  •  Nick Wenzel", size=15, color=WHITE, bold=True)
    textbox(s, 0.84, 5.93, 6.8, 0.3, "MWI Deep Learning Capstone  •  2 July 2026", size=12, color="94A3B8")
    textbox(s, 10.9, 6.85, 1.5, 0.25, "NILS • 0:20", size=9, color="94A3B8", align=PP_ALIGN.RIGHT)
    add_notes(s, "Nils — 20 seconds\nOpen with one sentence: ClaimGuard does not decide misconduct; it organizes evidence so a human can review it. Introduce both presenters and the three-part question: claim, reference, support.")

    # 2 — Motivation
    s = new_slide(prs, 2, "Why academic integrity checking matters", "Nils", "0:55", "Problem motivation")
    quote = box(s, 0.62, 1.35, 5.15, 4.85, fill=NAVY, line=NAVY)
    textbox(s, 0.95, 1.72, 4.45, 0.32, "A citation can be real…", size=21, color=CYAN, bold=True)
    textbox(s, 0.95, 2.28, 4.45, 1.22, "…and still not support the sentence that cites it.", size=29, color=WHITE, bold=True)
    textbox(s, 0.95, 4.08, 4.3, 0.8, "Manual review is slow; naive automation creates false accusations.", size=16, color="CBD5E1")
    pill(s, 0.95, 5.42, 2.38, "TRIAGE, NOT A VERDICT", AMBER, NAVY)
    add_bullets(
        s, 6.25, 1.58, 6.2, 3.9,
        [
            ("Scale", "Long reports contain dozens of checkable claims."),
            ("Subtlety", "Existence, relevance, and factual support are different questions."),
            ("Risk", "False positives can unfairly damage students or researchers."),
            ("Need", "Reviewers require traceable evidence, not a black-box score."),
        ], size=17, gap=15,
    )
    textbox(s, 6.28, 5.68, 5.8, 0.42, "Goal: make suspicious passages faster to inspect and easier to explain.", size=16, color=BLUE, bold=True)
    add_notes(s, "Nils — 55 seconds\nUse the contrast on the left. Explain that plagiarism, fabricated references, and unsupported citation use are distinct. ClaimGuard focuses on citation integrity and evidence support. Stress that the system is decision support, not an accusation engine.")

    # 3 — Research questions
    s = new_slide(prs, 3, "Three checks — one auditable report", "Nils", "0:45", "Problem framing")
    cards = [
        ("01", "Does this need a citation?", "Claim type, section, citation marker, severity", BLUE),
        ("02", "Does the source exist?", "IEEE parsing plus scholarly metadata APIs", PURPLE),
        ("03", "Does the source support it?", "Evidence retrieval plus five-way verification", GREEN),
    ]
    for i, (n, title, detail, color) in enumerate(cards):
        x = 0.72 + i * 4.18
        box(s, x, 1.55, 3.72, 3.82, fill=LIGHT, line=color)
        pill(s, x + 0.22, 1.8, 0.72, n, color)
        textbox(s, x + 0.22, 2.48, 3.2, 0.82, title, size=22, color=INK, bold=True)
        textbox(s, x + 0.22, 3.58, 3.15, 0.74, detail, size=15, color=MUTED)
    textbox(s, 1.05, 5.83, 11.2, 0.42, "Output: human-readable Markdown + complete JSON provenance for evaluation and audit", size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Nils — 45 seconds\nWalk left to right. These checks map directly to the three core modules. Mention that every decision remains traceable in JSON, while Markdown is the reviewer-facing view.")

    # 4 — Architecture
    s = new_slide(prs, 4, "Pipeline architecture", "Nils", "0:55", "System architecture")
    stages = [
        ("Ingest", "PDF/TXT extraction\nsection cleanup", BLUE),
        ("Claims", "type + citation need\ncontext linking", CYAN),
        ("References", "IEEE parsing\nmetadata validation", PURPLE),
        ("Evidence", "lexical / embeddings\nchunk retrieval", AMBER),
        ("Verdict", "heuristic / local /\nfrontier / LoRA", GREEN),
    ]
    for i, (title, detail, color) in enumerate(stages):
        x = 0.48 + i * 2.57
        add_architecture_stage(s, x, 1.58, 2.15, i + 1, title, detail, color)
        if i < 4:
            arrow = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 2.18), Inches(1.99), Inches(0.33), Inches(0.48))
            set_fill(arrow, LINE)
            set_line(arrow, LINE)
    # lower audit rail
    box(s, 0.74, 3.55, 11.84, 1.05, fill=NAVY, line=NAVY)
    textbox(s, 1.0, 3.77, 2.25, 0.35, "AUDIT RAIL", size=13, color=CYAN, bold=True)
    textbox(s, 3.05, 3.75, 8.9, 0.43, "evidence • confidence • model ID • latency • token usage • response/request IDs", size=15, color=WHITE)
    metric_card(s, 1.0, 5.05, 3.25, "Offline first", "Deterministic fallback", GREEN, "Network features are opt-in")
    metric_card(s, 5.03, 5.05, 3.25, "Same schema", "Comparable backends", BLUE, "Five support labels")
    metric_card(s, 9.05, 5.05, 3.25, "Human review", "Final decision", AMBER, "No automatic misconduct claim")
    add_notes(s, "Nils — 55 seconds\nFollow the five numbered boxes. Explain that references constrain evidence retrieval; the verifier receives bounded evidence rather than the whole internet. Point out the audit rail: this makes outputs traceable and comparable. Then introduce the module-coverage decision slide.")

    # 5 — Module coverage
    s = new_slide(prs, 5, "Module coverage: completed and deferred", "Nils", "0:55", "Rubric coverage")
    module_rows = [
        ("CORE 1–3", "Completed", "Claims/citations, bibliography, and RAG verification; required and evaluated.", "COMPLETE", GREEN),
        ("MODULE 5", "LoRA completed", "SciFact adapter trained, saved, and compared on aligned and transfer benchmarks.", "COMPLETE", BLUE),
        ("MODULE 4", "Completed via Fast-DetectGPT", "Official API evaluated on N=30; Binoculars' calibrated 14B pair exceeds 8 GB VRAM.", "COMPLETE", AMBER),
        ("MODULE 6", "Comparison scaffold only", "Shared import/evaluation code exists, but no scite/SemanticCite predictions were available.", "FUTURE", RED),
    ]
    for i, (tag, title, detail, status, color) in enumerate(module_rows):
        y = 1.34 + i * 1.13
        box(s, 0.7, y, 11.9, 0.98, fill=LIGHT if i % 2 == 0 else WHITE, line=LINE)
        pill(s, 0.9, y + 0.2, 1.42, tag, color)
        textbox(s, 2.55, y + 0.17, 2.85, 0.38, title, size=16, color=INK, bold=True)
        textbox(s, 5.42, y + 0.14, 5.05, 0.55, detail, size=12.5, color=MUTED)
        pill(s, 10.78, y + 0.32, 1.48, status, color)
    textbox(
        s,
        0.92,
        6.03,
        11.5,
        0.4,
        "Claimed as complete: Core Modules 1–3 + Optional Modules 4 and 5 • Module 6 remains future work.",
        size=14,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_notes(s, "Nils — 55 seconds\nThe three required core modules are complete and visibly evaluated. For optional depth, Module 5 is complete: we trained and saved the LoRA adapter and compared it on aligned and transfer benchmarks. Module 4 is fulfilled through Fast-DetectGPT. We considered Binoculars, but its calibrated Falcon pair totals about 14 billion parameters and does not fit the available 8 GB GPU; using a smaller pair would invalidate its published threshold. Module 6 has a working comparison/import scaffold, but without shared scite or SemanticCite predictions we do not claim a completed external-tool benchmark. This boundary is deliberate: implemented code is separated from measured evidence.")

    # 6 — Technical decisions
    s = new_slide(prs, 6, "Technical decisions and trade-offs", "Nick", "1:00", "Technical depth")
    rows = [
        ("Retrieval", "all-MiniLM-L6-v2", "Semantic matching; cosine fallback on Windows", BLUE),
        ("Metadata", "CrossRef • OpenAlex • S2", "Coverage across DOI, title, authors, OA status", PURPLE),
        ("Local model", "Ollama / qwen3:1.7b", "Private and offline; slower on this machine", GREEN),
        ("Frontier", "OpenAI gpt-5.4-mini", "Fast structured output; paid and external", CYAN),
        ("Fine-tuned", "SciFact LoRA", "Cheap local inference; three-label taxonomy", AMBER),
    ]
    for i, (area, choice, reason, color) in enumerate(rows):
        y = 1.42 + i * 0.98
        box(s, 0.7, y, 11.9, 0.76, fill=LIGHT if i % 2 == 0 else WHITE, line=LINE)
        pill(s, 0.88, y + 0.2, 1.22, area.upper(), color)
        textbox(s, 2.38, y + 0.16, 3.3, 0.4, choice, size=16, color=INK, bold=True)
        textbox(s, 5.82, y + 0.16, 6.35, 0.42, reason, size=14, color=MUTED)
    textbox(s, 0.9, 6.43, 11.6, 0.32, "Design principle: graceful degradation — APIs or embeddings may fail without stopping the report.", size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Nick — 60 seconds\nExplain why there is no single best backend. The APIs validate bibliography, not source support. The retriever can use embeddings but falls back to lexical search. Mention store=false for OpenAI and that local models keep text on-device.")

    # 7 — Demo
    s = new_slide(prs, 7, "Demo on our real RQ5 report", "Nick", "1:15", "Demo")
    box(s, 0.62, 1.35, 7.1, 1.35, fill=NAVY, line=NAVY)
    textbox(s, 0.9, 1.52, 6.55, 0.24, "POWERSHELL", size=10, color=CYAN, bold=True, font="Consolas")
    cmd = "python -m scripts.run_claimguard `\n  --input data/reports/NW2_RQ5_Report.pdf --verifier openai `\n  --output outputs/examples/rq5_openai.json --markdown-output outputs/examples/rq5_openai.md"
    textbox(s, 0.9, 1.82, 6.55, 0.68, cmd, size=11.5, color=WHITE, font="Consolas")
    metric_card(s, 8.05, 1.35, 1.35, "137", "sentences", BLUE)
    metric_card(s, 9.57, 1.35, 1.35, "65", "factual claims", PURPLE)
    metric_card(s, 11.09, 1.35, 1.35, "21", "citation flags", RED)
    textbox(s, 0.72, 3.08, 4.7, 0.35, "REAL CLAIM FROM SENTENCE 11", size=12, color=BLUE, bold=True)
    box(s, 0.65, 3.48, 5.65, 2.45, fill=LIGHT, line=LINE)
    textbox(s, 0.92, 3.78, 5.05, 0.95, "“RAG connects a language model to external knowledge without changing model weights.”", size=18, color=INK, bold=True)
    textbox(s, 0.92, 5.08, 5.0, 0.42, "Citations detected: [1, 3]  •  Top evidence retained", size=12, color=MUTED)
    # result cards
    box(s, 6.62, 3.48, 2.72, 2.45, fill="EFF6FF", line=BLUE)
    pill(s, 6.86, 3.72, 1.02, "OPENAI", BLUE)
    textbox(s, 6.86, 4.28, 2.2, 0.52, "insufficient\nevidence", size=18, color=BLUE, bold=True)
    textbox(s, 6.86, 5.17, 2.2, 0.35, "1.889 s • 506 tokens", size=11, color=MUTED)
    box(s, 9.62, 3.48, 2.72, 2.45, fill="ECFDF5", line=GREEN)
    pill(s, 9.86, 3.72, 1.02, "OLLAMA", GREEN)
    textbox(s, 9.86, 4.35, 2.2, 0.35, "supported", size=20, color=GREEN, bold=True)
    textbox(s, 9.86, 5.17, 2.2, 0.35, "40.902 s • local", size=11, color=MUTED)
    textbox(s, 0.82, 6.27, 11.65, 0.37, "Demo message: the evidence is visible — disagreement becomes reviewable instead of hidden.", size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Nick — 75 seconds\nIf live: run the prepared command, then open rq5_openai.md. If offline: use this slide. Read only the shortened claim. Point to citations and top evidence. Compare the two verdicts, then emphasize that disagreement is exactly why evidence and rationale must be shown.")

    # 8 — Results
    s = new_slide(prs, 8, "Key results: controlled tests vs. real reports", "Nick", "1:15", "Evaluation")
    chart_data = ChartData()
    chart_data.categories = ["Real RQ\nclaims", "Lexical\nverifier", "SciFact LoRA\naligned", "Fast-DetectGPT"]
    chart_data.add_series("F1 / Macro-F1", (0.437, 0.900, 0.367, 0.696))
    chart_shape = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.72), Inches(1.42), Inches(7.1), Inches(4.78), chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 1
    chart.value_axis.major_unit = 0.2
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = rgb(BLUE)
    chart.series[0].has_data_labels = True
    labels = chart.series[0].data_labels
    labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    labels.number_format = "0.000"
    labels.font.size = Pt(11)
    labels.font.bold = True
    metric_card(s, 8.25, 1.52, 3.92, "1.000", "Reference parse rate", GREEN, "30 authentic + controlled cases")
    metric_card(s, 8.25, 3.08, 3.92, "0.267", "Citation-needed F1", RED, "50 manually labeled RQ sentences")
    metric_card(s, 8.25, 4.64, 3.92, "0.924", "Fast-DetectGPT AUROC", AMBER, "15 human + 15 GPT-4 passages")
    textbox(s, 0.88, 6.37, 11.35, 0.34, "Main finding: synthetic benchmarks overestimate how well simple rules transfer to authentic PDFs.", size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Nick — 75 seconds\nStart with the 0.900 lexical verifier result, then contrast it with real claim classification at 0.437 Macro-F1. LoRA reached 0.367 on the label-aligned SciFact subset. Fast-DetectGPT reached 0.924 AUROC but only 0.533 recall, missing 7 of 15 generated passages. The takeaway is that every score needs its dataset and failure boundary.")

    # 9 — Model comparison
    s = new_slide(prs, 9, "Model comparison: speed, control, and taxonomy", "Nils", "1:15", "Local vs. frontier")
    headers = ["Backend", "Observed result", "Latency", "Best use", "Limitation"]
    widths = [1.65, 2.25, 1.45, 2.55, 3.55]
    xs = [0.68]
    for w in widths[:-1]:
        xs.append(xs[-1] + w)
    for i, (head, x, w) in enumerate(zip(headers, xs, widths)):
        box(s, x, 1.4, w, 0.55, fill=NAVY, line=WHITE, radius=False)
        textbox(s, x + 0.06, 1.53, w - 0.12, 0.25, head, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    rows = [
        ("Heuristic", "0.868 Macro-F1", "~0 ms", "Transparent baseline", "Relies on surface patterns", BLUE),
        ("LoRA", "0.172 Macro-F1", "6 ms", "Cheap local inference", "SciFact has 3 vs. 5 labels", AMBER),
        ("Ollama", "0.728 Macro-F1", "2.74 s", "Private local review", "Over-predicts full support", GREEN),
        ("OpenAI", "0.518 Macro-F1", "1.42 s", "Structured frontier review", "Paid/external; cost not logged", PURPLE),
    ]
    for r, row in enumerate(rows):
        y = 1.95 + r * 0.94
        values, color = row[:-1], row[-1]
        for c, (value, x, w) in enumerate(zip(values, xs, widths)):
            box(s, x, y, w, 0.82, fill=LIGHT if r % 2 == 0 else WHITE, line=LINE, radius=False)
            textbox(s, x + 0.08, y + 0.18, w - 0.16, 0.43, value, size=12.5 if c else 13.5, color=color if c == 0 else INK, bold=c == 0, align=PP_ALIGN.CENTER)
    box(s, 0.75, 5.95, 11.85, 0.68, fill="FFF7ED", line=AMBER)
    textbox(s, 0.98, 6.12, 11.35, 0.3, "Same 30 gold-labeled cases and identical retrieved evidence • 3 no-evidence abstentions", size=12.5, color="92400E", bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Nils — 75 seconds\nAll four backends saw the same retrieved evidence. The heuristic wins this small synthetic benchmark because cases reward direct lexical patterns. LoRA is fastest but cannot emit two of ClaimGuard's five labels. Ollama beat OpenAI here; OpenAI was faster. This is a benchmark-specific finding, not a universal model ranking.")

    # 10 — Works / limits
    s = new_slide(prs, 10, "What works — and where it breaks", "Nils", "0:55", "Limits")
    box(s, 0.68, 1.42, 5.8, 4.95, fill="ECFDF5", line=GREEN)
    pill(s, 0.96, 1.72, 1.48, "WORKS WELL", GREEN)
    add_bullets(s, 1.02, 2.35, 4.95, 3.25, [
        "End-to-end PDF → Markdown/JSON workflow",
        "Traceable evidence and backend metadata",
        "Robust IEEE parsing for Report 5",
        "LoRA and Fast-DetectGPT evaluated reproducibly",
    ], size=16, color=INK, gap=16)
    box(s, 6.84, 1.42, 5.8, 4.95, fill="FEF2F2", line=RED)
    pill(s, 7.12, 1.72, 1.52, "LIMITATIONS", RED)
    add_bullets(s, 7.18, 2.35, 4.95, 3.25, [
        "PDF extraction introduces spacing/heading noise",
        "Real claim detection generalizes poorly",
        "Repeated S2/OpenAlex HTTP 429 in manual runs",
        "Fast-DetectGPT missed 7/15 generated passages",
    ], size=16, color=INK, gap=16)
    textbox(s, 1.25, 6.55, 10.85, 0.3, "Ethical boundary: a flag prioritizes human review; it never proves misconduct or AI authorship.", size=14.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Nils — 55 seconds\nGive two concrete successes and two failures. In manual runs on 7 July 2026, Semantic Scholar requests without an API key and OpenAlex requests with a configured mailto parameter repeatedly returned HTTP 429, even after waiting several minutes and changing the input file. Present this as an observed service condition, not a systematic API-availability benchmark; it can reduce abstract/full-text evidence and lead to insufficient-evidence results. Fast-DetectGPT had no human false positives in this small set, but missed 7 of 15 generated passages. The official benchmark has no citation markers, so no AI-score/citation correlation is claimed. Private reports were not sent to the external detector.")

    # 11 — Conclusion
    s = new_slide(prs, 11, "Takeaways", "Nick + Nils", "0:35", "Conclusion")
    set_fill(s.background, NAVY)
    # mask standard header area with dark overlay and restyle
    overlay = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    set_fill(overlay, NAVY)
    set_line(overlay, NAVY)
    textbox(s, 0.8, 0.58, 4.0, 0.35, "CONCLUSION", size=11, color=CYAN, bold=True)
    textbox(s, 0.8, 1.1, 11.5, 0.7, "ClaimGuard makes integrity review auditable.", size=34, color=WHITE, bold=True)
    takeaways = [
        ("01", "Pipeline", "Claims, references, evidence, and verdicts in one workflow."),
        ("02", "Evidence", "Real-report evaluation exposes where synthetic scores mislead."),
        ("03", "Responsibility", "Use automation to focus human attention — not replace judgment."),
    ]
    for i, (n, title, detail) in enumerate(takeaways):
        x = 0.8 + i * 4.1
        box(s, x, 2.35, 3.65, 2.28, fill=NAVY_2, line="334155")
        pill(s, x + 0.22, 2.61, 0.7, n, BLUE)
        textbox(s, x + 0.22, 3.15, 3.1, 0.35, title, size=18, color=CYAN, bold=True)
        textbox(s, x + 0.22, 3.64, 3.12, 0.62, detail, size=13, color="CBD5E1")
    textbox(s, 0.82, 5.55, 11.5, 0.45, "Next: double annotation, external-tool benchmark, and consented citation-pattern study.", size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 0.82, 6.75, 5.0, 0.25, "Nick closes • Nils transitions to Q&A", size=9, color="94A3B8")
    add_notes(s, "Nick — 20 seconds; Nils — 15 seconds\nNick summarizes pipeline and evidence. Nils closes with the ethical boundary and invites questions. Transition immediately to the Q&A slide.")

    # 12 — Q&A
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_fill(s.background, NAVY)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.73), Inches(0.88), Inches(1.82), Inches(1.82))
    set_fill(circle, BLUE)
    set_line(circle, CYAN, 2)
    textbox(s, 5.73, 1.28, 1.82, 0.72, "?", size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 1.2, 3.05, 10.9, 0.78, "Questions & discussion", size=37, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 2.15, 4.05, 9.0, 0.45, "5 minutes", size=20, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 2.0, 5.2, 9.3, 0.55, "Architecture  •  Evaluation  •  Model trade-offs  •  Ethics", size=15, color="CBD5E1", align=PP_ALIGN.CENTER)
    textbox(s, 0.75, 6.82, 4.5, 0.25, "Nils Wagner • Nick Wenzel", size=10, color="94A3B8")
    textbox(s, 10.8, 6.82, 1.7, 0.25, "05:00 Q&A", size=10, color="94A3B8", align=PP_ALIGN.RIGHT)
    add_notes(s, "Both — 5 minutes\nSuggested ownership: Nils answers architecture/evaluation; Nick answers model/API/demo questions. If asked about the low real-claim score, lead with the measured domain shift and single-annotator limitation.")

    return prs


def write_notes(prs: Presentation) -> None:
    lines = ["# ClaimGuard presentation notes", "", "10-minute talk + 5-minute Q&A.", ""]
    for index, slide in enumerate(prs.slides, start=1):
        notes = slide.notes_slide.notes_text_frame.text.strip()
        lines.extend([f"## Slide {index}", "", notes, ""])
    NOTES_OUT.write_text("\n".join(lines), encoding="utf-8")


def validate(prs: Presentation) -> None:
    if len(prs.slides) != 12:
        raise RuntimeError(f"Expected 12 slides, got {len(prs.slides)}")
    for index, slide in enumerate(prs.slides, start=1):
        if not slide.notes_slide.notes_text_frame.text.strip():
            raise RuntimeError(f"Slide {index} has no speaker notes")
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                raise RuntimeError(f"Slide {index} has a shape outside the canvas")
            if shape.left + shape.width > W + Inches(0.02):
                raise RuntimeError(f"Slide {index} has a shape beyond the right edge")
            if shape.top + shape.height > H + Inches(0.02):
                raise RuntimeError(f"Slide {index} has a shape beyond the bottom edge")


def main() -> None:
    prs = build_deck()
    validate(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    write_notes(prs)
    print(f"Wrote {OUT}")
    print(f"Wrote {NOTES_OUT}")


if __name__ == "__main__":
    main()
