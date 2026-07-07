"""Create a fully German ClaimGuard deck from the editable English presentation."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation


HERE = Path(__file__).resolve().parent
BUILD_DIR = HERE / "build"
SOURCE = BUILD_DIR / "ClaimGuard_Capstone_Presentation.pptx"
OUTPUT = BUILD_DIR / "ClaimGuard_Capstone_Praesentation_DE.pptx"
NOTES_OUTPUT = BUILD_DIR / "ClaimGuard_Praesentation_Notizen_DE.md"


TRANSLATIONS = {
    "15 MIN • TEAM PRESENTATION": "15 MIN • TEAMPRÄSENTATION",
    "Evidence-backed academic integrity checking": "Evidenzbasierte Prüfung akademischer Integrität",
    "From PDF claims to reference validation and source-support verdicts": "Von Aussagen im PDF zur Referenzprüfung und Quellenbewertung",
    "MWI Deep Learning Capstone  •  2 July 2026": "MWI Deep Learning Capstone  •  3. Juli 2026",
    "PROBLEM MOTIVATION": "PROBLEMMOTIVATION",
    "Why academic integrity checking matters": "Warum die Prüfung akademischer Integrität wichtig ist",
    "A citation can be real…": "Eine Quelle kann echt sein …",
    "…and still not support the sentence that cites it.": "… und trotzdem die zitierte Aussage nicht stützen.",
    "Manual review is slow; naive automation creates false accusations.": "Manuelle Prüfung ist langsam; naive Automatisierung erzeugt falsche Verdächtigungen.",
    "TRIAGE, NOT A VERDICT": "VORPRÜFUNG, KEIN URTEIL",
    "•  Scale  Long reports contain dozens of checkable claims.": "•  Umfang  Lange Berichte enthalten viele überprüfbare Aussagen.",
    "•  Subtlety  Existence, relevance, and factual support are different questions.": "•  Feinheit  Existenz, Relevanz und faktische Stützung sind verschiedene Fragen.",
    "•  Risk  False positives can unfairly damage students or researchers.": "•  Risiko  Fehlalarme können Studierenden und Forschenden unfair schaden.",
    "•  Need  Reviewers require traceable evidence, not a black-box score.": "•  Bedarf  Prüfende brauchen nachvollziehbare Evidenz statt einer Blackbox-Zahl.",
    "Goal: make suspicious passages faster to inspect and easier to explain.": "Ziel: Auffällige Passagen schneller prüfen und verständlich begründen.",
    "PROBLEM FRAMING": "PROBLEMSTELLUNG",
    "Three checks — one auditable report": "Drei Prüfungen — ein nachvollziehbarer Bericht",
    "Does this need a citation?": "Benötigt das eine Quelle?",
    "Claim type, section, citation marker, severity": "Aussagetyp, Abschnitt, Zitiermarker, Schweregrad",
    "Does the source exist?": "Existiert die Quelle?",
    "IEEE parsing plus scholarly metadata APIs": "IEEE-Parsing plus wissenschaftliche Metadaten-APIs",
    "Does the source support it?": "Stützt die Quelle die Aussage?",
    "Evidence retrieval plus five-way verification": "Evidenzsuche plus fünfstufige Verifikation",
    "Output: human-readable Markdown + complete JSON provenance for evaluation and audit": "Ausgabe: lesbares Markdown + vollständige JSON-Provenienz für Evaluation und Audit",
    "SYSTEM ARCHITECTURE": "SYSTEMARCHITEKTUR",
    "Pipeline architecture": "Pipeline-Architektur",
    "Ingest": "Einlesen",
    "PDF/TXT extraction\vsection cleanup": "PDF/TXT-Extraktion\vAbschnittsbereinigung",
    "Claims": "Aussagen",
    "type + citation need\vcontext linking": "Typ + Quellenbedarf\vKontextverknüpfung",
    "References": "Referenzen",
    "IEEE parsing\vmetadata validation": "IEEE-Parsing\vMetadatenprüfung",
    "Evidence": "Evidenz",
    "lexical / embeddings\vchunk retrieval": "lexikalisch / Embeddings\vChunk-Suche",
    "Verdict": "Bewertung",
    "heuristic / local /\vfrontier / LoRA": "Heuristik / lokal /\vFrontier / LoRA",
    "AUDIT RAIL": "AUDIT-SPUR",
    "evidence • confidence • model ID • latency • token usage • response/request IDs": "Evidenz • Konfidenz • Modell-ID • Latenz • Tokenverbrauch • Antwort-/Anfrage-IDs",
    "Offline first": "Offline zuerst",
    "Deterministic fallback": "Regelbasierter Fallback",
    "Network features are opt-in": "Netzwerkfunktionen nur optional",
    "Same schema": "Ein Schema",
    "Comparable backends": "Vergleichbare Backends",
    "Five support labels": "Fünf Unterstützungslabels",
    "Human review": "Mensch prüft",
    "Final decision": "Endgültige Entscheidung",
    "No automatic misconduct claim": "Kein automatischer Täuschungsvorwurf",
    "RUBRIC COVERAGE": "MODULABDECKUNG",
    "Module coverage: implemented, adapted, deferred": "Modulabdeckung: umgesetzt, angepasst, zurückgestellt",
    "Module coverage: completed and deferred": "Modulabdeckung: erfüllt und offen",
    "CORE 1–3": "KERN 1–3",
    "Completed": "Abgeschlossen",
    "Claims/citations, bibliography, and RAG verification; required and evaluated.": "Claims/Zitate, Bibliografie und RAG-Verifikation; verpflichtend und evaluiert.",
    "COMPLETE": "ERFÜLLT",
    "MODULE 5": "MODUL 5",
    "LoRA completed": "LoRA abgeschlossen",
    "SciFact adapter trained, saved, and compared on aligned and transfer benchmarks.": "SciFact-Adapter trainiert, gespeichert und labelgleich sowie im Transfer verglichen.",
    "MODULE 4": "MODUL 4",
    "Completed via Fast-DetectGPT": "Über Fast-DetectGPT erfüllt",
    "Official API evaluated on N=30; Binoculars' calibrated 14B pair exceeds 8 GB VRAM.": "Offizielle API auf N=30 evaluiert; Binoculars' kalibriertes 14B-Paar übersteigt 8 GB VRAM.",
    "ADAPTED": "ANGEPASST",
    "MODULE 6": "MODUL 6",
    "Comparison scaffold only": "Nur Vergleichsgerüst",
    "Shared import/evaluation code exists, but no scite/SemanticCite predictions were available.": "Import- und Evaluationscode existiert, aber keine scite-/SemanticCite-Vorhersagen lagen vor.",
    "FUTURE": "OFFEN",
    "Claimed as complete: Core Modules 1–3 + Optional Modules 4 and 5 • Module 6 remains future work.": "Als erfüllt ausgewiesen: Kernmodule 1–3 + optionale Module 4 und 5 • Modul 6 bleibt Future Work.",
    "TECHNICAL DEPTH": "TECHNISCHE TIEFE",
    "Technical decisions and trade-offs": "Technische Entscheidungen und Zielkonflikte",
    "RETRIEVAL": "EVIDENZSUCHE",
    "Semantic matching; cosine fallback on Windows": "Semantischer Abgleich; Kosinus-Fallback unter Windows",
    "METADATA": "METADATEN",
    "Coverage across DOI, title, authors, OA status": "Abdeckung von DOI, Titel, Autoren und Open-Access-Status",
    "LOCAL MODEL": "LOKAL",
    "Private and offline; slower on this machine": "Privat und offline; auf diesem Rechner langsamer",
    "FRONTIER": "FRONTIER",
    "Fast structured output; paid and external": "Schnelle strukturierte Ausgabe; kostenpflichtig und extern",
    "FINE-TUNED": "FINE-TUNING",
    "Cheap local inference; three-label taxonomy": "Günstige lokale Inferenz; Taxonomie mit drei Labels",
    "Design principle: graceful degradation — APIs or embeddings may fail without stopping the report.": "Designprinzip: kontrollierte Degradation — API- oder Embedding-Fehler stoppen den Bericht nicht.",
    "Demo on our real RQ5 report": "Demo mit unserem echten RQ5-Bericht",
    "sentences": "Sätze",
    "factual claims": "faktische Aussagen",
    "citation flags": "Quellenhinweise",
    "REAL CLAIM FROM SENTENCE 11": "ECHTE AUSSAGE AUS SATZ 11 (ÜBERSETZT)",
    "“RAG connects a language model to external knowledge without changing model weights.”": "„RAG verbindet ein Sprachmodell mit externem Wissen, ohne Modellgewichte zu verändern.“",
    "Citations detected: [1, 3]  •  Top evidence retained": "Erkannte Quellen: [1, 3]  •  Beste Evidenz gespeichert",
    "insufficient\vevidence": "unzureichende\vEvidenz",
    "supported": "gestützt",
    "40.902 s • local": "40,902 s • lokal",
    "Demo message: the evidence is visible — disagreement becomes reviewable instead of hidden.": "Kernaussage: Die Evidenz ist sichtbar — Unterschiede werden prüfbar statt verborgen.",
    "EVALUATION": "EVALUATION",
    "Key results: controlled tests vs. real reports": "Zentrale Ergebnisse: kontrollierte Tests vs. echte Berichte",
    "IEEE parser accuracy": "Genauigkeit des IEEE-Parsers",
    "5-field-complete test cases": "Fünf Testfälle mit vollständigen Feldern",
    "Citation-needed F1": "F1 für Quellenbedarf",
    "50 manually labeled RQ sentences": "50 manuell annotierte RQ-Sätze",
    "SciFact validation Macro-F1": "SciFact-Validierung: Macro-F1",
    "450 validation pairs": "450 Validierungspaare",
    "Main finding: synthetic benchmarks overestimate how well simple rules transfer to authentic PDFs.": "Hauptergebnis: Synthetische Benchmarks überschätzen den Transfer einfacher Regeln auf echte PDFs.",
    "Reference parse rate": "Referenz-Parse-Rate",
    "30 authentic + controlled cases": "30 echte + kontrollierte Fälle",
    "Fast-DetectGPT AUROC": "Fast-DetectGPT AUROC",
    "15 human + 15 GPT-4 passages": "15 menschliche + 15 GPT-4-Passagen",
    "LOCAL VS. FRONTIER": "LOKAL VS. FRONTIER",
    "Model comparison: speed, control, and taxonomy": "Modellvergleich: Geschwindigkeit, Kontrolle und Taxonomie",
    "Backend": "Backend",
    "Observed result": "Beobachtetes Ergebnis",
    "Latency": "Latenz",
    "Best use": "Geeigneter Einsatz",
    "Limitation": "Einschränkung",
    "Heuristic": "Heuristik",
    "Transparent baseline": "Transparente Baseline",
    "Relies on surface patterns": "Beruht auf Oberflächenmustern",
    "Cheap local inference": "Günstige lokale Inferenz",
    "SciFact has 3 vs. 5 labels": "SciFact hat 3 statt 5 Labels",
    "Over-predicts full support": "Überschätzt vollständige Stützung",
    "Paid/external; cost not logged": "Kostenpflichtig/extern; Kosten nicht erfasst",
    "Same 30 gold-labeled cases and identical retrieved evidence • 3 no-evidence abstentions": "Dieselben 30 Gold-Fälle und identische Evidenz • 3 Enthaltungen ohne Evidenz",
    "Private local review": "Private lokale Prüfung",
    "Structured frontier review": "Strukturierte Frontier-Prüfung",
    "LIMITS": "GRENZEN",
    "What works — and where it breaks": "Was funktioniert — und wo es scheitert",
    "WORKS WELL": "FUNKTIONIERT GUT",
    "•  End-to-end PDF → Markdown/JSON workflow": "•  Durchgängiger Ablauf von PDF zu Markdown/JSON",
    "•  Traceable evidence and backend metadata": "•  Nachvollziehbare Evidenz und Backend-Metadaten",
    "•  Robust IEEE parsing for Report 5": "•  Robustes IEEE-Parsing für Bericht 5",
    "•  Offline fallback and real CUDA LoRA path": "•  Offline-Fallback und echter CUDA-LoRA-Pfad",
    "•  LoRA and Fast-DetectGPT evaluated reproducibly": "•  LoRA und Fast-DetectGPT reproduzierbar evaluiert",
    "LIMITATIONS": "EINSCHRÄNKUNGEN",
    "•  PDF extraction introduces spacing/heading noise": "•  PDF-Extraktion erzeugt Abstands- und Überschriftenfehler",
    "•  Real claim detection generalizes poorly": "•  Erkennung echter Aussagen generalisiert schwach",
    "•  Repeated S2/OpenAlex HTTP 429 in manual runs": "•  Wiederholte S2/OpenAlex-HTTP-429-Fehler in manuellen Läufen",
    "•  Binoculars' calibrated 14B pair exceeds 8 GB VRAM": "•  Binoculars' kalibriertes 14B-Paar übersteigt 8 GB VRAM",
    "•  Fast-DetectGPT missed 7/15 generated passages": "•  Fast-DetectGPT übersah 7/15 generierte Passagen",
    "Ethical boundary: a flag prioritizes human review; it never proves misconduct or AI authorship.": "Ethische Grenze: Ein Hinweis priorisiert menschliche Prüfung; er beweist weder Täuschung noch KI-Autorschaft.",
    "CONCLUSION": "FAZIT",
    "Takeaways": "Kernaussagen",
    "ClaimGuard makes integrity review auditable.": "ClaimGuard macht Integritätsprüfungen nachvollziehbar.",
    "Pipeline": "Pipeline",
    "Claims, references, evidence, and verdicts in one workflow.": "Aussagen, Referenzen, Evidenz und Bewertungen in einem Ablauf.",
    "Real-report evaluation exposes where synthetic scores mislead.": "Evaluation mit echten Berichten zeigt, wo synthetische Werte täuschen.",
    "Responsibility": "Verantwortung",
    "Use automation to focus human attention — not replace judgment.": "Automatisierung lenkt menschliche Aufmerksamkeit — sie ersetzt kein Urteil.",
    "Next: expand double-annotated real claims and run the full OpenAI/Ollama gold benchmark.": "Nächster Schritt: doppelt annotierte Echtdaten erweitern und den vollständigen OpenAI/Ollama-Benchmark ausführen.",
    "Next: double annotation, external-tool benchmark, and consented citation-pattern study.": "Nächster Schritt: Doppelannotation, externer Tool-Benchmark und consentierte Zitiermusterstudie.",
    "Nick closes • Nils transitions to Q&A": "Nick fasst zusammen • Nils leitet zur Fragerunde über",
    "Questions & discussion": "Fragen & Diskussion",
    "5 minutes": "5 Minuten",
    "Architecture  •  Evaluation  •  Model trade-offs  •  Ethics": "Architektur  •  Evaluation  •  Modell-Zielkonflikte  •  Ethik",
    "05:00 Q&A": "05:00 FRAGEN",
}


CHART_TRANSLATIONS = {
    "Real RQ\nclaims": "Echte RQ-\nAussagen",
    "Lexical\nverifier": "Lexikalische\nPrüfung",
    "SciFact LoRA\naligned": "SciFact LoRA\nlabelgleich",
    "Fast-DetectGPT": "Fast-DetectGPT",
    "F1 / Macro-F1": "F1 / Macro-F1",
}


GERMAN_NOTES = [
    "Nils — 20 Sekunden\nMit einem Satz eröffnen: ClaimGuard entscheidet nicht über Täuschung, sondern ordnet Evidenz für eine menschliche Prüfung. Beide Vortragenden sowie die drei Fragen Aussage, Referenz und Stützung vorstellen.",
    "Nils — 55 Sekunden\nDen Kontrast links nutzen. Plagiat, erfundene Referenzen und nicht gestützte Zitate sind unterschiedliche Probleme. ClaimGuard konzentriert sich auf Zitierintegrität und Evidenz. Das System unterstützt Entscheidungen; es erhebt keinen Vorwurf.",
    "Nils — 45 Sekunden\nVon links nach rechts erklären. Die drei Fragen entsprechen den drei Kernmodulen. Jede Entscheidung bleibt im JSON nachvollziehbar; Markdown ist die Ansicht für Prüfende.",
    "Nils — 55 Sekunden\nDen fünf nummerierten Schritten folgen. Referenzen begrenzen die Evidenzsuche; der Verifikator erhält begrenzte Evidenz statt des gesamten Internets. Die Audit-Spur macht Ausgaben nachvollziehbar und vergleichbar. Danach die Entscheidungen zur Modulabdeckung erklären.",
    "Nils — 55 Sekunden\nDie drei verpflichtenden Kernmodule sind vollständig umgesetzt und evaluiert. Als optionale Vertiefung ist Modul 5 abgeschlossen: Der LoRA-Adapter wurde trainiert, gespeichert und auf labelgleichen sowie Transfer-Benchmarks verglichen. Modul 4 ist über Fast-DetectGPT erfüllt. Binoculars wurde geprüft, aber das kalibrierte Falcon-Paar umfasst etwa 14 Milliarden Parameter und passt nicht in die verfügbare 8-GB-GPU; ein kleineres Paar würde den publizierten Schwellenwert ungültig machen. Für Modul 6 existiert das Vergleichs- und Importgerüst, aber ohne gemeinsame scite- oder SemanticCite-Vorhersagen wird kein fertiger externer Benchmark behauptet. Diese Grenze trennt implementierten Code bewusst von tatsächlich gemessener Evidenz.",
    "Nick — 60 Sekunden\nEs gibt kein einzelnes bestes Backend. Die APIs prüfen bibliografische Daten, nicht die inhaltliche Stützung. Die Evidenzsuche nutzt Embeddings und fällt bei Bedarf auf lexikalische Suche zurück. OpenAI nutzt store=false; lokale Modelle halten Texte auf dem Gerät.",
    "Nick — 75 Sekunden\nBeim Live-Durchlauf den echten RQ5-Bericht und danach den Modellvergleich zeigen. Heuristik und Ollama erhalten dieselben Claims und dieselbe Evidenz. Abweichungen sind kein automatischer Fehlernachweis, sondern gezielte Stellen für die menschliche Prüfung. Falls die Demo nicht läuft, die statische Folie oder ClaimGuard_UI_Demo.mp4 verwenden.",
    "Nick — 75 Sekunden\nMit 0,900 Macro-F1 der lexikalischen Prüfung beginnen und mit 0,437 bei echten Aussagen kontrastieren. LoRA erreicht auf dem labelgleichen SciFact-Subset 0,367. Fast-DetectGPT erreicht 0,924 AUROC, aber nur 0,533 Recall und übersieht 7 von 15 generierten Passagen. Jeder Wert braucht Datensatz und Fehlergrenze.",
    "Nils — 75 Sekunden\nAlle vier Backends sahen dieselbe Evidenz. Die Heuristik gewinnt diesen kleinen synthetischen Benchmark, weil direkte lexikalische Muster belohnt werden. LoRA ist am schnellsten, kann aber zwei der fünf Labels nicht ausgeben. Ollama schlägt OpenAI hier; OpenAI ist schneller. Das ist kein universelles Modellranking.",
    "Nils — 55 Sekunden\nJe zwei konkrete Stärken und Schwächen nennen. In manuellen Läufen am 7. Juli 2026 lieferten Semantic-Scholar-Anfragen ohne API-Key und OpenAlex-Anfragen mit konfiguriertem mailto-Parameter wiederholt HTTP 429, auch nach mehreren Minuten Wartezeit und mit einer neuen Eingabedatei. Das ist eine beobachtete Dienstbedingung und kein systematischer API-Verfügbarkeitsbenchmark; fehlende Abstract- oder Volltextevidenz kann dadurch zu unzureichender Evidenz führen. Fast-DetectGPT erzeugte in diesem kleinen Set keinen menschlichen Fehlalarm, übersah aber 7 von 15 generierten Texten. Der offizielle Datensatz enthält keine Zitiermarker; eine Korrelation wird daher nicht behauptet. Private Berichte wurden nicht an den externen Detektor gesendet.",
    "Nick — 20 Sekunden; Nils — 15 Sekunden\nNick fasst Pipeline und Evidenz zusammen. Nils schließt mit der ethischen Grenze und lädt zu Fragen ein.",
    "Beide — 5 Minuten\nVorgeschlagene Aufteilung: Nils beantwortet Fragen zu Architektur und Evaluation; Nick übernimmt Modelle, APIs und Demo. Bei Fragen zum niedrigen Echtdatenwert zuerst auf Domain Shift und die Einzelannotation eingehen.",
]


def replace_paragraph_text(paragraph, new_text: str) -> None:
    """Replace text while keeping the first run's formatting."""
    if paragraph.runs:
        parts = new_text.split("\v")
        paragraph.runs[0].text = parts[0]
        for run in paragraph.runs[1:]:
            run.text = ""
        for part in parts[1:]:
            paragraph.add_line_break()
            paragraph.add_run().text = part
    else:
        paragraph.text = new_text.replace("\v", "\n")


def translate_visible_text(prs: Presentation) -> set[str]:
    untranslated: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                old = paragraph.text
                if old in TRANSLATIONS:
                    replace_paragraph_text(paragraph, TRANSLATIONS[old])
                elif old.strip() and any(ch.isalpha() for ch in old):
                    untranslated.add(old)
    return untranslated


def translate_chart(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_chart", False):
                continue
            for value in shape.chart._chartSpace.xpath(".//c:v"):
                if value.text in CHART_TRANSLATIONS:
                    value.text = CHART_TRANSLATIONS[value.text]


def write_notes(prs: Presentation) -> None:
    lines = ["# ClaimGuard – Präsentationsnotizen", "", "10 Minuten Vortrag + 5 Minuten Fragen.", ""]
    for index, (slide, note) in enumerate(zip(prs.slides, GERMAN_NOTES, strict=True), start=1):
        slide.notes_slide.notes_text_frame.text = note
        lines.extend([f"## Folie {index}", "", note, ""])
    NOTES_OUTPUT.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)
    prs = Presentation(SOURCE)
    if len(prs.slides) != len(GERMAN_NOTES):
        raise RuntimeError(f"Erwartet: {len(GERMAN_NOTES)} Folien; gefunden: {len(prs.slides)}")
    untranslated = translate_visible_text(prs)
    translate_chart(prs)
    write_notes(prs)
    prs.core_properties.title = "ClaimGuard — Evidenzbasierte Prüfung akademischer Integrität"
    prs.core_properties.subject = "Deutsche MWI-Deep-Learning-Capstone-Präsentation"
    prs.save(OUTPUT)
    print(f"Erstellt: {OUTPUT}")
    print(f"Erstellt: {NOTES_OUTPUT}")
    print("Bewusst unverändert (Namen, Befehle, Modelle und Messwerte):")
    for text in sorted(untranslated):
        print(f"  {text!r}")


if __name__ == "__main__":
    main()
